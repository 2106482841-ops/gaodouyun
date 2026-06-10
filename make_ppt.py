#!/usr/bin/env python3
"""Generate 高斗云招生方案 PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C_PRIMARY  = RGBColor(0x1A, 0x3C, 0x6E)
C_ACCENT   = RGBColor(0x2D, 0x7D, 0xD0)
C_ACCENT2  = RGBColor(0xE8, 0x6C, 0x00)
C_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK     = RGBColor(0x33, 0x33, 0x33)
C_GRAY     = RGBColor(0x66, 0x66, 0x66)
C_LIGHT    = RGBColor(0xAA, 0xBB, 0xCC)
C_GREEN    = RGBColor(0x27, 0xAE, 0x60)
C_RED      = RGBColor(0xE7, 0x4C, 0x3C)
C_YELLOW   = RGBColor(0xF3, 0x9C, 0x12)
C_PURPLE   = RGBColor(0x8E, 0x44, 0xAD)

FN = "Microsoft YaHei"

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, rounded=False, line=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background() if not line else None
    s.shadow.inherit = False
    return s

def oval(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def tb(slide, l, t, w, h, text, sz=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FN
    p.alignment = align
    return box

def bullets(slide, l, t, w, h, items, sz=16, color=C_DARK):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FN
        p.space_after = Pt(4)
    return box

def line_bar(slide, l, t, w, color=C_ACCENT, thick=Pt(4)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, thick)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def header(slide, text):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_PRIMARY)
    tb(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
       text, sz=28, bold=True, color=C_WHITE)

def footer_bar(slide, text, color=C_PRIMARY):
    rect(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), color)
    tb(slide, Inches(0.8), Inches(6.85), Inches(11), Inches(0.5),
       text, sz=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ======== Slide 1: Cover ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_PRIMARY)
rect(sl, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_ACCENT)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.08), C_ACCENT2)
tb(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
   "高斗云单词速记工具", sz=48, bold=True, color=C_WHITE)
line_bar(sl, Inches(1.5), Inches(2.8), Inches(4), C_ACCENT2, Pt(5))
tb(sl, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
   "2026年暑假招生方案", sz=36, color=RGBColor(0xBB,0xCC,0xDD))
tb(sl, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
   "右脑图像记忆法 + 六脉神剑 + 阅读宝 AI 诊断", sz=20, color=RGBColor(0x88,0xAA,0xCC))
tb(sl, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
   "AAA 级信用企业认证  |  1 个月投资回报周期  |  双业务引擎", sz=16, color=C_LIGHT)

# ======== Slide 2: Market Background ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_WHITE)
header(sl, "新教改后，英语学习的底层逻辑已经变了")

blocks = [
    ("政策导向: 从应试到应用", "高考英语 150 分主科地位不变\n考核从语法知识转向真实语境运用\n强调跨文化沟通与实际能力", C_ACCENT),
    ("教材难度飙升", "八年级新版 Unit1: 52词升至75词\n高中/初三词汇提前下放\n八年级成为初中英语分水岭", C_ACCENT2),
    ("词汇要求远超考纲", "新课标: 中考2000词、高考3500词\n但应对真实语篇实际需更多\n不只是会拼写，更要会使用", C_GREEN),
]
for i, (title, desc, clr) in enumerate(blocks):
    x = Inches(0.8 + i * 4.2)
    y = Inches(1.8)
    rect(sl, x, y, Inches(3.8), Inches(4.5), C_LIGHT_BG, rounded=True)
    rect(sl, x, y, Inches(3.8), Inches(0.08), clr)
    tb(sl, x+Inches(0.3), y+Inches(0.3), Inches(3.2), Inches(0.6),
       title, sz=18, bold=True, color=clr)
    tb(sl, x+Inches(0.3), y+Inches(1.2), Inches(3.2), Inches(3.0),
       desc, sz=16, color=C_DARK)

tb(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
   "核心矛盾: 从是否努力转变为方法是否高效", sz=18, bold=True,
   color=C_RED, align=PP_ALIGN.CENTER)

# ======== Slide 3: Pain Points ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_WHITE)
header(sl, "传统记忆法的三大困境")

pains = [
    ("记忆效率低下", "日均1-2小时死记硬背\n24小时遗忘率高达67%\n陷入背了忘、忘了背的死循环"),
    ("认词不等于会用", "能认单词但在实际场景无法运用\n阅读长难句、写作表达卡壳\n知识储备无法转化为能力"),
    ("学习兴趣丧失", "单调枯燥导致心理疲劳和厌学\n海量词汇的无力感加频繁遗忘\n自信心崩塌，越学越怕"),
]
for i, (title, desc) in enumerate(pains):
    x = Inches(0.8 + i * 4.2)
    y = Inches(1.8)
    rect(sl, x, y, Inches(3.8), Inches(4.5), RGBColor(0xFD,0xED,0xEC), rounded=True)
    rect(sl, x, y, Inches(3.8), Inches(0.08), C_RED)
    tb(sl, x+Inches(0.3), y+Inches(0.3), Inches(3.2), Inches(0.5),
       title, sz=20, bold=True, color=C_RED)
    tb(sl, x+Inches(0.3), y+Inches(1.0), Inches(3.2), Inches(3.0),
       desc, sz=16, color=C_DARK)

tb(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
   "学生真正缺的不是背单词的时间，而是一套科学高效的方法",
   sz=18, bold=True, color=C_RED, align=PP_ALIGN.CENTER)

# ======== Slide 4: Competitors ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_WHITE)
header(sl, "竞争格局: 巨头林立，但都有明显短板")

competitors = [
    ("百词斩", "~31%", "图片联想记忆", "认图不认词\n脱离图片后词义回忆模糊\n长期记忆留存率低", C_LIGHT_BG),
    ("扇贝单词", "~26%", "词库分级+词根词缀", "内容偏学术化\n低龄学生趣味性不足\n碎片化学习者吸引力弱", RGBColor(0xE8,0xF8,0xF5)),
    ("墨墨背单词", "~19%", "艾宾浩斯自适应算法", "界面传统枯燥\n缺乏游戏化激励\n用户长期动力易衰减", RGBColor(0xFE,0xF9,0xE7)),
    ("高斗云", "新兴", "右脑图像+六脉神剑", "方法论本质革新\n避开同质化竞争\nB端模式差异化突围", RGBColor(0xE8,0xF0,0xFE)),
]
for i, (name, share, method, weak, bg) in enumerate(competitors):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.8)
    rect(sl, x, y, Inches(2.9), Inches(4.8), bg, rounded=True)
    tb(sl, x+Inches(0.2), y+Inches(0.3), Inches(2.5), Inches(0.5),
       name, sz=22, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.2), y+Inches(0.8), Inches(2.5), Inches(0.4),
       share, sz=14, color=C_GRAY, align=PP_ALIGN.CENTER)
    line_bar(sl, x+Inches(0.4), y+Inches(1.2), Inches(2.1), C_ACCENT, Pt(2))
    tb(sl, x+Inches(0.2), y+Inches(1.4), Inches(2.5), Inches(0.6),
       "核心: " + method, sz=13, bold=True, color=C_ACCENT)
    tb(sl, x+Inches(0.2), y+Inches(2.2), Inches(2.5), Inches(2.0),
       weak, sz=14, color=C_DARK)

footer_bar(sl, "高斗云机会: 以科学方法论 + B端代理模式，避开C端红海，精准切入学校与教培机构场景")

# ======== Slide 5: Brain Science ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_WHITE)
header(sl, "右脑图像记忆法: 为什么更高效？")

# Left brain
rect(sl, Inches(0.5), Inches(1.6), Inches(5.8), Inches(4.0), RGBColor(0xFD,0xED,0xEC), rounded=True)
rect(sl, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.08), C_RED)
tb(sl, Inches(0.8), Inches(1.9), Inches(5.2), Inches(0.5),
   "左脑: 文字记忆（传统方式）", sz=20, bold=True, color=C_RED)
bullets(sl, Inches(0.8), Inches(2.5), Inches(5.2), Inches(2.8), [
    "线性串行处理器，逐字解码",
    "容量仅约 7+2 组块，极易遗忘",
    "24小时遗忘率高达 67%",
    "单次仅承载 5.2 个单元",
    "消耗大量认知资源，容易疲劳",
], sz=15, color=C_DARK)

# Right brain
rect(sl, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.0), RGBColor(0xE8,0xF8,0xF5), rounded=True)
rect(sl, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.08), C_GREEN)
tb(sl, Inches(7.3), Inches(1.9), Inches(5.2), Inches(0.5),
   "右脑: 图像记忆（高斗云方法）", sz=20, bold=True, color=C_GREEN)
bullets(sl, Inches(7.3), Inches(2.5), Inches(5.2), Inches(2.8), [
    "多维并行处理器，全息快照存储",
    "近乎无限的记忆容量",
    "24小时遗忘率仅 34%",
    "单次可承载 15.8 个单元（3倍以上）",
    "认知负荷降低 57%",
], sz=15, color=C_DARK)

tb(sl, Inches(0.5), Inches(5.8), Inches(11.5), Inches(0.4),
   "数据来源: Brady et al., 2008; Sharot et al., 2007; MIT大脑与认知科学实验室, 2022",
   sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)

rect(sl, Inches(3.0), Inches(6.3), Inches(7.3), Inches(0.7), C_PRIMARY, rounded=True)
tb(sl, Inches(3.0), Inches(6.35), Inches(7.3), Inches(0.5),
   "科学实证: 图像记忆效率是文字记忆的 3 倍以上",
   sz=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ======== Slide 6: Six Methods ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_WHITE)
header(sl, "独创六脉神剑记忆法")

methods = [
    ("谐音法", "ambulance -> 俺不能死\n发音转中文谐音建立关联\n一秒速记，有趣难忘", C_RED),
    ("形象法", "eye -> 形似眼睛轮廓\n字母组合对应具象图形\n右脑图像替代死记硬背", C_ACCENT2),
    ("词根词缀法", "un- + happy = unhappy\n掌握构词逻辑一通百通\n像拼积木一样扩充词汇", C_GREEN),
    ("熟词法", "class + room = classroom\n已知单词推导新单词含义\n用旧知带动新知", C_ACCENT),
    ("成图法", "字母含义绘制成画\n想象力构建趣味场景\n创意记忆，难以遗忘", C_PURPLE),
    ("智能系统", "AI 自动匹配错题复习\n个性化推送学习内容\n科学规划遗忘曲线", RGBColor(0x00,0x96,0x88)),
]
for i, (name, desc, clr) in enumerate(methods):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(1.8 + row * 2.8)
    rect(sl, x, y, Inches(3.8), Inches(2.4), C_LIGHT_BG, rounded=True)
    rect(sl, x, y, Inches(3.8), Inches(0.08), clr)
    tb(sl, x+Inches(0.3), y+Inches(0.2), Inches(3.2), Inches(0.5),
       name, sz=20, bold=True, color=clr)
    tb(sl, x+Inches(0.3), y+Inches(0.7), Inches(3.2), Inches(1.5),
       desc, sz=14, color=C_DARK)

footer_bar(sl, "好记 - 好全 - 好听 - 好玩 - 好管: 五大产品价值")

# ======== Slide 7: 阅读宝 ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_WHITE)
header(sl, "阅读宝 --- AI 赋能个性化英语提分利器")

features = [
    ("AI 精准诊断", "五大维度全面分析\n词汇、短语、阅读、翻译、语法\n精准抓取薄弱环节"),
    ("定制专属学案", "根据诊断结果量身打造\n匹配真实水平推送学习内容\n针对性解决薄弱点"),
    ("科学系统训练", "五大专项学习方案\n词汇、阅读、语法、短语、翻译\n动态优化循序渐进"),
    ("全程可视化", "实时反馈学情数据\n进步变化一目了然\n家长随时掌握学习动态"),
]
for i, (title, desc) in enumerate(features):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.0)
    rect(sl, x, y, Inches(2.9), Inches(4.5), C_LIGHT_BG, rounded=True)
    rect(sl, x, y, Inches(2.9), Inches(0.08), C_ACCENT)
    tb(sl, x+Inches(0.2), y+Inches(0.3), Inches(2.5), Inches(0.5),
       title, sz=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.2), y+Inches(1.0), Inches(2.5), Inches(3.2),
       desc, sz=14, color=C_DARK)

footer_bar(sl, "三种模式: 引流获客（免费诊断）| 课内赋能（AI+课堂）| 增值增收（专项提分营）")

# ======== Slide 8: Summer Campaign ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_WHITE)
header(sl, "暑假招生三板斧")

campaigns = [
    ("9.9 单词体验课", "低价引流", "邀约家长线下体验\n孩子体验3-6个单词即\n接受右脑图像记忆法\n家长认可自然买单", C_RED),
    ("99 四次速记课", "转化成交", "叶乘风校长手把手代教\n全套话术框架直接复制\n从引流到成交全程\n可复制的执行方案", C_ACCENT2),
    ("单词挑战赛", "地推获客", "校门口和校区附近摆摊\n1分钟挑战3个单词\n挑战成功赢奖品或抽红包\n游戏化吸引加收集线索", C_GREEN),
]
for i, (title, subtitle, desc, clr) in enumerate(campaigns):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.8)
    rect(sl, x, y, Inches(3.8), Inches(4.8), C_LIGHT_BG, rounded=True)
    rect(sl, x, y, Inches(3.8), Inches(0.08), clr)
    tb(sl, x+Inches(0.3), y+Inches(0.3), Inches(3.2), Inches(0.5),
       title, sz=22, bold=True, color=clr, align=PP_ALIGN.CENTER)
    tag = rect(sl, x+Inches(1.0), y+Inches(0.9), Inches(1.8), Inches(0.4), clr, rounded=True)
    tb(sl, x+Inches(1.0), y+Inches(0.92), Inches(1.8), Inches(0.35),
       subtitle, sz=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.3), y+Inches(1.6), Inches(3.2), Inches(2.8),
       desc, sz=15, color=C_DARK)

footer_bar(sl, "三环联动: 低价引流 -> 低价转化 -> 地推获客，形成暑假招生闭环")

# ======== Slide 9: ROI ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_PRIMARY)
rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0F,0x2A,0x4E))
tb(sl, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
   "投资回报测算 --- 极低门槛，极高爆发力", sz=28, bold=True, color=C_WHITE)

cards = [
    ("单词训练系统", "294,000", "元/年", "500学员 x 588元单生利润", C_ACCENT),
    ("阅读宝", "480,000", "元/年", "100学员 x 4800元/年", C_GREEN),
    ("线下陪跑", "15,000+", "元/期", "单班半年，10人班\n老师成本100-200/节", C_ACCENT2),
]
for i, (title, amt, unit, note, clr) in enumerate(cards):
    x = Inches(0.5 + i * 4.3)
    y = Inches(1.8)
    rect(sl, x, y, Inches(3.9), Inches(3.0), C_WHITE, rounded=True)
    rect(sl, x, y, Inches(3.9), Inches(0.08), clr)
    tb(sl, x+Inches(0.3), y+Inches(0.3), Inches(3.3), Inches(0.4), title, sz=16, color=C_GRAY)
    tb(sl, x+Inches(0.3), y+Inches(0.8), Inches(2.5), Inches(0.7),
       "\u00a5" + amt, sz=40, bold=True, color=clr)
    tb(sl, x+Inches(2.8), y+Inches(1.2), Inches(0.8), Inches(0.4), unit, sz=14, color=C_GRAY)
    tb(sl, x+Inches(0.3), y+Inches(1.9), Inches(3.3), Inches(0.8), note, sz=13, color=C_DARK)

# Summary
rect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), C_WHITE, rounded=True)
rect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.08), C_ACCENT2)
tb(sl, Inches(1.0), Inches(5.4), Inches(5), Inches(0.6), "单店年度总利润目标", sz=18, color=C_GRAY)
tb(sl, Inches(1.0), Inches(5.9), Inches(3), Inches(0.8),
   "\u00a5774,000", sz=44, bold=True, color=C_RED)
tb(sl, Inches(5.5), Inches(5.5), Inches(6), Inches(0.5),
   "极速投资回报周期: 1 个月", sz=20, bold=True, color=C_ACCENT2)
tb(sl, Inches(5.5), Inches(6.0), Inches(6), Inches(0.5),
   "双业务引擎叠加 | 轻资产运营 | 科技教育蓝海", sz=16, color=C_DARK)

# ======== Slide 10: Closing ========
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, C_PRIMARY)
rect(sl, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_ACCENT)
rect(sl, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), C_ACCENT2)
tb(sl, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
   "携手高斗云 共赢新未来", sz=44, bold=True, color=C_WHITE)
line_bar(sl, Inches(1.5), Inches(3.2), Inches(4), C_ACCENT2, Pt(5))
tb(sl, Inches(1.5), Inches(3.6), Inches(10), Inches(1.5),
   "期待与您携手，共创辉煌！\n\nAAA级信用企业认证\n权威机构认证，彰显品牌实力与信誉\n是值得您信赖的长期战略合作伙伴",
   sz=20, color=RGBColor(0xBB,0xCC,0xDD))
tb(sl, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
   "高斗云单词速记工具  |  2026年暑假招生方案", sz=14, color=C_LIGHT)

out = "/Users/lvjinna/.openclaw/workspace/高斗云招生方案.pptx"
prs.save(out)
print(f"OK! Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
